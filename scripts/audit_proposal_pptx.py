#!/usr/bin/env python3
"""
audit_proposal_pptx.py — objective delivery QA for proposal-ppt skill.

The skill's Hard Rules forbid claiming a PowerPoint deck was delivered when it
was not, and forbid shipping decks that still carry placeholder text. This
script turns those rules from "trust the agent" into "verify the file".

It runs on pure standard library (zipfile + xml.etree) so it works anywhere.
If `python-pptx` is installed, deeper checks (blank frame detection, per-slide
text density) are added automatically.

Usage:
    python3 audit_proposal_pptx.py <deck.pptx> [--script presenter.md] [--strict]

Exit code:
    0  no errors (warnings may still be present)
    1  one or more errors found, or the file is not a valid deck

Heavily inspired by the sibling `audit_static_site.py` pattern in
business-website-skill: independent of the agent, errors/warnings/info graded.
"""

from __future__ import annotations

import argparse
import re
import sys
import xml.etree.ElementTree as ET
import zipfile
from pathlib import Path

# Markers the skill itself tells the agent to use for unknown data, plus the
# usual leftover-placeholder suspects. If any of these survive into the final
# deck, that is an error unless the page is explicitly a "known unknowns" page.
PLACEHOLDER_MARKERS = [
    "待补充", "待确认", "暂无公开数据", "需客户确认",
    "TODO", "TO DO", "FIXME", "TBD", "TBC", "XXX",
    "lorem ipsum", "placeholder", "占位", "示例文字",
    "[插入", "[填写", "[待",
]

# Per the skill's output contract, a usable presenter script is expected to
# have roughly one block per slide. These heading patterns cover the bilingual
# formats the skill emits.
SCRIPT_SLIDE_PATTERNS = [
    re.compile(r"^#{1,6}\s*(?:slide|幻灯片|页|第\s*\d+\s*页|p\d+|page\s*\d+)", re.I),
    re.compile(r"^###\s*\d+[\.、:：]", re.I),
]

CONTENT_TYPES = "[Content_Types].xml"
SLIDE_RE = re.compile(r"^ppt/slides/slide\d+\.xml$")
NOTES_RE = re.compile(r"^ppt/notesSlides/notesSlide\d+\.xml$")

# OOXML namespaces used for text extraction.
NS = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}


class Report:
    def __init__(self):
        self.errors = []
        self.warnings = []
        self.info = []

    def err(self, msg):
        self.errors.append(msg)

    def warn(self, msg):
        self.warnings.append(msg)

    def note(self, msg):
        self.info.append(msg)

    def print(self, stream=sys.stdout):
        for m in self.info:
            stream.write(f"  [INFO ] {m}\n")
        for m in self.warnings:
            stream.write(f"  [WARN ] {m}\n")
        for m in self.errors:
            stream.write(f"  [ERROR] {m}\n")


def extract_text(xml_bytes: bytes) -> str:
    """Return concatenated <a:t> text from a slide/notes XML part."""
    try:
        root = ET.fromstring(xml_bytes)
    except ET.ParseError:
        return ""
    parts = [node.text for node in root.iter() if node.tag.endswith("}t") and node.text]
    return "\n".join(parts)


def count_script_slides(script_path: Path) -> int | None:
    """Best-effort count of slide blocks in the presenter .md script."""
    if not script_path or not script_path.exists():
        return None
    count = 0
    for line in script_path.read_text(encoding="utf-8", errors="replace").splitlines():
        if any(p.match(line.strip()) for p in SCRIPT_SLIDE_PATTERNS):
            count += 1
    return count or None


def detect_placeholder_text(text: str, page_label: str, report: Report):
    """Flag placeholder markers that leaked into shipped text."""
    for marker in PLACEHOLDER_MARKERS:
        if marker.lower() in text.lower():
            report.err(
                f"{page_label}: placeholder marker '{marker}' found in shipped text "
                f"— replace with real content or mark the page as a known-unknowns page."
            )


def python_pptx_blank_frames(deck_path: Path, slide_count: int, report: Report, template_mode: bool = False):
    """When python-pptx is available, detect oversized empty shapes, text-box
    overflow risk, and element-to-element overlap."""
    try:
        from pptx import Presentation  # type: ignore
        from pptx.util import Emu
    except Exception:
        report.note("python-pptx not installed — skipping blank-frame and density checks.")
        return
    try:
        prs = Presentation(str(deck_path))
    except Exception as exc:  # corrupt or not a real deck
        report.err(f"python-pptx could not open the deck: {exc}")
        return

    slide_w = prs.slide_width or 1
    slide_h = prs.slide_height or 1
    for idx, slide in enumerate(prs.slides, start=1):
        non_empty = 0
        blank_large = 0
        text_boxes = []  # (left, top, right, bottom, text) in EMU for overlap checks
        for shape in slide.shapes:
            area_ratio = (shape.width or 0) * (shape.height or 0) / (slide_w * slide_h)
            has_text = bool(shape.has_text_frame and shape.text_frame.text.strip())
            has_image = shape.shape_type == 13  # PICTURE
            if has_text or has_image:
                non_empty += 1
            elif area_ratio > 0.25:
                # A large shape with no text and no image = likely a dead frame.
                blank_large += 1
            # Text-box overflow risk: estimate whether the text fits its box.
            # This is a heuristic — real wrapping depends on the host renderer —
            # but it catches the common case of a long CJK string crammed into a
            # small auto-sized title box, which is the #1 source of ragged wraps.
            # Skipped in template mode: blank-template boxes are intentionally
            # small and meant to be filled, so sizing hints there are noise.
            if has_text and shape.width and shape.height and not template_mode:
                text = shape.text_frame.text
                # Approx EMU-per-CJK-glyph from the largest run font size; fall
                # back to a default body size if sizes are not set.
                size_pt = 18
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size:
                            size_pt = run.font.size.pt
                            break
                    if any(run.font.size for run in para.runs):
                        break
                emu_per_pt = 12700  # 1 pt = 12700 EMU
                box_w_pt = shape.width / emu_per_pt
                # Roughly 1 CJK glyph ≈ font size in pt of width.
                chars_per_line = max(1, box_w_pt / size_pt)
                lines_in_text = sum(max(1, (len(p.text) + chars_per_line - 1) // chars_per_line)
                                    for p in shape.text_frame.paragraphs if p.text)
                box_h_lines = max(1, (shape.height / emu_per_pt) / (size_pt * 1.2))
                if lines_in_text > box_h_lines + 1:
                    snippet = text.replace("\n", " ")[:24]
                    report.warn(
                        f"slide {idx}: text box may overflow — ~{lines_in_text} line(s) of text "
                        f"in a box that fits ~{int(box_h_lines)} line(s) ('{snippet}…'). "
                        f"Tighten copy or enlarge the box; verify in the renderer."
                    )
                # Collect for overlap check (expand nothing; use shape box).
                left = shape.left or 0
                top = shape.top or 0
                text_boxes.append((left, top, left + (shape.width or 0), top + (shape.height or 0),
                                   text.replace("\n", " ")[:18]))
        if blank_large and non_empty <= 1:
            report.warn(
                f"slide {idx}: {blank_large} large empty frame(s) with little content "
                f"— possible dead placeholder region."
            )
        if non_empty == 0:
            report.warn(f"slide {idx}: no text or image content detected.")
        # Element-to-element overlap check (text boxes only).
        margin = int(slide_w * 0.01)  # 1% of slide width tolerance
        for i in range(len(text_boxes)):
            for j in range(i + 1, len(text_boxes)):
                a = text_boxes[i]
                b = text_boxes[j]
                ox = min(a[2], b[2]) - max(a[0], b[0])
                oy = min(a[3], b[3]) - max(a[1], b[1])
                if ox > margin and oy > margin:
                    report.err(
                        f"slide {idx}: two text boxes overlap by {int(ox/12700)}×{int(oy/12700)} pt "
                        f"— ('{a[4]}…' ↔ '{b[4]}…'). Overlap reads as a broken deck; treat as a hard fail."
                    )
    if len(prs.slides) != slide_count:
        report.warn(
            f"python-pptx sees {len(prs.slides)} slides but the package lists {slide_count} "
            f"— package structure may be inconsistent."
        )


def audit(deck_path: Path, script_path: Path | None, report: Report, template_mode: bool = False):
    if not deck_path.exists():
        report.err(f"file not found: {deck_path}")
        return

    # 1. ZIP / Office Open XML validity.
    if not zipfile.is_zipfile(deck_path):
        report.err("not a valid .pptx (not a zip archive).")
        return
    try:
        zf = zipfile.ZipFile(deck_path)
    except zipfile.BadZipFile as exc:
        report.err(f"zip could not be opened: {exc}")
        return

    names = zf.namelist()
    if CONTENT_TYPES not in names:
        report.err("missing [Content_Types].xml — not a valid Office Open XML package.")
        return
    if not any(n.startswith("ppt/") for n in names):
        report.err("no ppt/ parts found — this is not a PowerPoint package.")
        return

    slide_names = sorted(n for n in names if SLIDE_RE.match(n))
    notes_names = sorted(n for n in names if NOTES_RE.match(n))
    slide_count = len(slide_names)

    if slide_count == 0:
        report.err("no slides found in ppt/slides/ — deck is empty.")
    else:
        report.note(f"slide count: {slide_count}")
        if slide_count < 5:
            report.warn(f"only {slide_count} slide(s) — most commercial proposals need more.")

    # 2. Per-slide text + placeholder scan (skipped for blank templates).
    if template_mode:
        report.note("template mode: placeholder-text scan skipped (fill-in markers allowed).")
    else:
        for name in slide_names:
            text = extract_text(zf.read(name))
            detect_placeholder_text(text, name, report)

    # 3. Notes / presenter script coverage.
    report.note(f"notes slides: {len(notes_names)}")
    if slide_count and len(notes_names) < slide_count:
        report.warn(
            f"only {len(notes_names)} notes slide(s) for {slide_count} slide(s) "
            f"— the skill expects a presenter talk track per slide."
        )

    # 4. Cross-check against the external .md script, if provided.
    if script_path:
        md_count = count_script_slides(script_path)
        if md_count is None:
            report.warn(f"could not detect slide blocks in script: {script_path}")
        elif slide_count and md_count != slide_count:
            report.warn(
                f"script has {md_count} slide block(s) but deck has {slide_count} slide(s) "
                f"— they should stay aligned."
            )
        else:
            report.note(f"script slide blocks ({md_count}) match deck slide count.")
    elif slide_count:
        report.warn("no presenter script passed with --script; cannot verify talk-track alignment.")

    zf.close()

    # 5. Deeper checks when python-pptx is available.
    if slide_count:
        python_pptx_blank_frames(deck_path, slide_count, report, template_mode=template_mode)


def main(argv=None):
    parser = argparse.ArgumentParser(
        description="Objective QA for proposal-ppt deliverables.",
    )
    parser.add_argument("deck", type=Path, help="path to the .pptx deck")
    parser.add_argument("--script", type=Path, default=None,
                        help="optional presenter script .md to cross-check slide count")
    parser.add_argument("--strict", action="store_true",
                        help="treat warnings as errors (non-zero exit on any warning)")
    parser.add_argument("--template", action="store_true",
                        help="audit a blank fallback template: keep structural, slide-count, "
                             "notes, and blank-frame checks, but skip the placeholder-text scan "
                             "(templates are expected to carry fill-in markers).")
    args = parser.parse_args(argv)

    report = Report()
    audit(args.deck, args.script, report, template_mode=args.template)

    deck_label = args.deck.name
    print(f"\n=== audit: {deck_label} ===")
    if not (report.errors or report.warnings or report.info):
        print("  no issues found.")
    else:
        report.print()
    print("=== result: ", end="")
    if report.errors:
        print(f"FAIL ({len(report.errors)} error(s), {len(report.warnings)} warning(s))")
    elif report.warnings and args.strict:
        print(f"FAIL in strict mode ({len(report.warnings)} warning(s))")
    elif report.warnings:
        print(f"PASS with {len(report.warnings)} warning(s)")
    else:
        print("PASS")

    has_errors = bool(report.errors) or (args.strict and report.warnings)
    return 1 if has_errors else 0


if __name__ == "__main__":
    sys.exit(main())
